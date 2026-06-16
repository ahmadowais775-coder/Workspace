from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Color scheme
DARK_BLUE = RGBColor(0x1B, 0x2A, 0x4A)
MEDIUM_BLUE = RGBColor(0x2C, 0x3E, 0x6B)
LIGHT_BLUE = RGBColor(0x3A, 0x7B, 0xD5)
ACCENT_BLUE = RGBColor(0x5D, 0xA0, 0xE8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF0, 0xF4, 0xF8)
ORANGE = RGBColor(0xF3, 0x9C, 0x12)


def add_background(slide, color):
    """Set slide background color."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_title_slide(prs):
    """Slide 1: Title Slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    add_background(slide, DARK_BLUE)

    # Add decorative shape
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(9), Inches(-1), Inches(5), Inches(5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = MEDIUM_BLUE
    shape.line.fill.background()

    # Title
    txBox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(2))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Artificial Intelligence"
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # Subtitle
    p2 = tf.add_paragraph()
    p2.text = "An Overview of AI Technology"
    p2.font.size = Pt(24)
    p2.font.color.rgb = ACCENT_BLUE
    p2.space_before = Pt(20)

    # Author/Date
    txBox2 = slide.shapes.add_textbox(Inches(1), Inches(5), Inches(8), Inches(1))
    tf2 = txBox2.text_frame
    p3 = tf2.paragraphs[0]
    p3.text = "June 2026 | Technology Presentation"
    p3.font.size = Pt(16)
    p3.font.color.rgb = RGBColor(0xAA, 0xBB, 0xCC)


def add_agenda_slide(prs):
    """Slide 2: Agenda / Table of Contents"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, WHITE)

    # Side bar
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.4), Inches(7.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = LIGHT_BLUE
    shape.line.fill.background()

    # Title
    txBox = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(10), Inches(1))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Agenda"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE

    # Agenda items
    items = [
        "1. What is Artificial Intelligence?",
        "2. History of AI",
        "3. Types of AI",
        "4. Key Applications",
        "5. Machine Learning & Deep Learning",
        "6. AI in Everyday Life",
        "7. Challenges & Ethics",
        "8. The Future of AI",
    ]

    txBox2 = slide.shapes.add_textbox(Inches(1.5), Inches(1.8), Inches(10), Inches(5))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf2.paragraphs[0]
        else:
            p = tf2.add_paragraph()
        p.text = item
        p.font.size = Pt(20)
        p.font.color.rgb = MEDIUM_BLUE
        p.space_before = Pt(14)


def add_what_is_ai_slide(prs):
    """Slide 3: What is AI?"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, LIGHT_GRAY)

    # Title
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11), Inches(1))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "What is Artificial Intelligence?"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE

    # Definition box
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.8), Inches(11.5), Inches(2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = LIGHT_BLUE
    shape.line.width = Pt(2)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.3)
    tf.margin_top = Inches(0.3)
    p = tf.paragraphs[0]
    p.text = "Artificial Intelligence (AI) is the simulation of human intelligence by machines"
    p.font.size = Pt(20)
    p.font.color.rgb = DARK_BLUE
    p2 = tf.add_paragraph()
    p2.text = "that are programmed to think, learn, and make decisions."
    p2.font.size = Pt(20)
    p2.font.color.rgb = DARK_BLUE

    # Key points
    points = [
        "Mimics cognitive functions like learning and problem-solving",
        "Can process vast amounts of data rapidly",
        "Improves over time through experience (Machine Learning)",
        "Ranges from narrow/specific tasks to general intelligence",
    ]

    txBox2 = slide.shapes.add_textbox(Inches(1.2), Inches(4.2), Inches(10), Inches(3))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    for i, point in enumerate(points):
        if i == 0:
            p = tf2.paragraphs[0]
        else:
            p = tf2.add_paragraph()
        p.text = f"\u2022  {point}"
        p.font.size = Pt(18)
        p.font.color.rgb = MEDIUM_BLUE
        p.space_before = Pt(12)


def add_history_slide(prs):
    """Slide 4: History of AI"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, WHITE)

    # Side bar
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.4), Inches(7.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ORANGE
    shape.line.fill.background()

    # Title
    txBox = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(11), Inches(1))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "History of AI"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE

    # Timeline items
    timeline = [
        ("1950", "Alan Turing proposes the Turing Test"),
        ("1956", "Term 'Artificial Intelligence' coined at Dartmouth Conference"),
        ("1966", "ELIZA chatbot created at MIT"),
        ("1997", "IBM Deep Blue defeats chess champion Garry Kasparov"),
        ("2011", "IBM Watson wins Jeopardy!"),
        ("2016", "Google DeepMind's AlphaGo defeats Go champion"),
        ("2022+", "Large Language Models (ChatGPT, etc.) transform AI landscape"),
    ]

    y_pos = 1.7
    for year, event in timeline:
        # Year badge
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), Inches(y_pos), Inches(1.5), Inches(0.5))
        shape.fill.solid()
        shape.fill.fore_color.rgb = LIGHT_BLUE
        shape.line.fill.background()
        tf = shape.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.text = year
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

        # Event text
        txBox = slide.shapes.add_textbox(Inches(3), Inches(y_pos), Inches(9), Inches(0.5))
        tf2 = txBox.text_frame
        tf2.vertical_anchor = MSO_ANCHOR.MIDDLE
        p2 = tf2.paragraphs[0]
        p2.text = event
        p2.font.size = Pt(16)
        p2.font.color.rgb = DARK_BLUE

        y_pos += 0.78


def add_types_of_ai_slide(prs):
    """Slide 5: Types of AI"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, DARK_BLUE)

    # Title
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11), Inches(1))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Types of Artificial Intelligence"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # Three columns
    types_data = [
        ("Narrow AI (ANI)", "Designed for a specific\ntask. Most AI today\nfalls into this category.\n\nExamples:\n- Siri / Alexa\n- Spam filters\n- Recommendation systems"),
        ("General AI (AGI)", "Hypothetical AI with\nhuman-level intelligence\nacross all domains.\n\nCapabilities:\n- Reasoning\n- Learning\n- Understanding context"),
        ("Super AI (ASI)", "Theoretical AI that\nsurpasses human\nintelligence entirely.\n\nCharacteristics:\n- Self-aware\n- Creative\n- Emotionally intelligent"),
    ]

    x_pos = 0.5
    for title, desc in types_data:
        # Card
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x_pos), Inches(1.8), Inches(3.8), Inches(5))
        shape.fill.solid()
        shape.fill.fore_color.rgb = MEDIUM_BLUE
        shape.line.fill.background()

        # Card title
        txBox = slide.shapes.add_textbox(Inches(x_pos + 0.3), Inches(2.0), Inches(3.4), Inches(0.8))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = ACCENT_BLUE
        p.alignment = PP_ALIGN.CENTER

        # Card description
        txBox2 = slide.shapes.add_textbox(Inches(x_pos + 0.3), Inches(2.8), Inches(3.4), Inches(3.8))
        tf2 = txBox2.text_frame
        tf2.word_wrap = True
        p2 = tf2.paragraphs[0]
        p2.text = desc
        p2.font.size = Pt(14)
        p2.font.color.rgb = WHITE

        x_pos += 4.2


def add_applications_slide(prs):
    """Slide 6: Key Applications"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, LIGHT_GRAY)

    # Title
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11), Inches(1))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Key Applications of AI"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE

    # Applications grid (2x3)
    apps = [
        ("Healthcare", "Disease diagnosis, drug\ndiscovery, medical imaging"),
        ("Finance", "Fraud detection, algorithmic\ntrading, risk assessment"),
        ("Transportation", "Self-driving cars, traffic\noptimization, logistics"),
        ("Education", "Personalized learning,\nautomated grading, tutoring"),
        ("Entertainment", "Content recommendations,\ngame AI, content creation"),
        ("Manufacturing", "Quality control, predictive\nmaintenance, robotics"),
    ]

    positions = [
        (0.5, 1.8), (4.5, 1.8), (8.5, 1.8),
        (0.5, 4.5), (4.5, 4.5), (8.5, 4.5),
    ]

    for (title, desc), (x, y) in zip(apps, positions):
        # Card
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(3.8), Inches(2.2))
        shape.fill.solid()
        shape.fill.fore_color.rgb = WHITE
        shape.line.color.rgb = LIGHT_BLUE
        shape.line.width = Pt(1.5)

        # Title
        txBox = slide.shapes.add_textbox(Inches(x + 0.2), Inches(y + 0.2), Inches(3.4), Inches(0.6))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = LIGHT_BLUE

        # Description
        txBox2 = slide.shapes.add_textbox(Inches(x + 0.2), Inches(y + 0.9), Inches(3.4), Inches(1.2))
        tf2 = txBox2.text_frame
        tf2.word_wrap = True
        p2 = tf2.paragraphs[0]
        p2.text = desc
        p2.font.size = Pt(14)
        p2.font.color.rgb = DARK_BLUE


def add_ml_dl_slide(prs):
    """Slide 7: Machine Learning & Deep Learning"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, WHITE)

    # Side bar
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.4), Inches(7.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = LIGHT_BLUE
    shape.line.fill.background()

    # Title
    txBox = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(11), Inches(1))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Machine Learning & Deep Learning"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE

    # ML Section
    shape1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.8), Inches(5.8), Inches(5))
    shape1.fill.solid()
    shape1.fill.fore_color.rgb = LIGHT_GRAY
    shape1.line.fill.background()

    txBox = slide.shapes.add_textbox(Inches(1.1), Inches(2.0), Inches(5.4), Inches(4.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Machine Learning"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = LIGHT_BLUE

    ml_points = [
        "Subset of AI that learns from data",
        "Types: Supervised, Unsupervised, Reinforcement",
        "Uses algorithms to identify patterns",
        "Improves with more data and experience",
    ]
    for point in ml_points:
        p = tf.add_paragraph()
        p.text = f"\u2022  {point}"
        p.font.size = Pt(15)
        p.font.color.rgb = DARK_BLUE
        p.space_before = Pt(10)

    # DL Section
    shape2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7), Inches(1.8), Inches(5.8), Inches(5))
    shape2.fill.solid()
    shape2.fill.fore_color.rgb = LIGHT_GRAY
    shape2.line.fill.background()

    txBox2 = slide.shapes.add_textbox(Inches(7.3), Inches(2.0), Inches(5.4), Inches(4.5))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = "Deep Learning"
    p2.font.size = Pt(22)
    p2.font.bold = True
    p2.font.color.rgb = ORANGE

    dl_points = [
        "Subset of ML using neural networks",
        "Inspired by human brain structure",
        "Excels at: image/speech recognition",
        "Powers: GPT, DALL-E, self-driving cars",
    ]
    for point in dl_points:
        p = tf2.add_paragraph()
        p.text = f"\u2022  {point}"
        p.font.size = Pt(15)
        p.font.color.rgb = DARK_BLUE
        p.space_before = Pt(10)


def add_challenges_slide(prs):
    """Slide 8: Challenges & Ethics"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, DARK_BLUE)

    # Title
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11), Inches(1))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Challenges & Ethical Considerations"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # Challenge items
    challenges = [
        ("Bias & Fairness", "AI systems can perpetuate or amplify existing biases in training data"),
        ("Privacy", "Large-scale data collection raises serious privacy concerns"),
        ("Job Displacement", "Automation may replace certain jobs, requiring workforce adaptation"),
        ("Transparency", "\"Black box\" AI models are difficult to interpret and explain"),
        ("Security", "AI systems can be vulnerable to adversarial attacks"),
        ("Regulation", "Need for governance frameworks to ensure responsible AI development"),
    ]

    y_pos = 1.7
    for title, desc in challenges:
        # Bullet shape
        shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1), Inches(y_pos + 0.15), Inches(0.3), Inches(0.3))
        shape.fill.solid()
        shape.fill.fore_color.rgb = ORANGE
        shape.line.fill.background()

        # Title
        txBox = slide.shapes.add_textbox(Inches(1.7), Inches(y_pos - 0.05), Inches(4), Inches(0.6))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = ACCENT_BLUE

        # Description
        txBox2 = slide.shapes.add_textbox(Inches(5.5), Inches(y_pos - 0.05), Inches(7), Inches(0.6))
        tf2 = txBox2.text_frame
        tf2.word_wrap = True
        p2 = tf2.paragraphs[0]
        p2.text = desc
        p2.font.size = Pt(15)
        p2.font.color.rgb = WHITE

        y_pos += 0.92


def add_future_slide(prs):
    """Slide 9: The Future of AI"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, LIGHT_GRAY)

    # Title
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11), Inches(1))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "The Future of AI"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE

    # Future predictions
    future_items = [
        "AI-powered personalized medicine and healthcare",
        "Autonomous vehicles becoming mainstream",
        "AI assistants handling complex professional tasks",
        "Breakthroughs in scientific research accelerated by AI",
        "Human-AI collaboration as standard in workplaces",
        "Advances toward Artificial General Intelligence (AGI)",
        "Stricter AI governance and ethical frameworks globally",
    ]

    txBox2 = slide.shapes.add_textbox(Inches(1.2), Inches(1.8), Inches(11), Inches(5))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    for i, item in enumerate(future_items):
        if i == 0:
            p = tf2.paragraphs[0]
        else:
            p = tf2.add_paragraph()
        p.text = f"\u27A4  {item}"
        p.font.size = Pt(20)
        p.font.color.rgb = MEDIUM_BLUE
        p.space_before = Pt(16)


def add_thank_you_slide(prs):
    """Slide 10: Thank You"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, DARK_BLUE)

    # Decorative circle
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-2), Inches(4), Inches(6), Inches(6))
    shape.fill.solid()
    shape.fill.fore_color.rgb = MEDIUM_BLUE
    shape.line.fill.background()

    # Thank you text
    txBox = slide.shapes.add_textbox(Inches(2), Inches(2.5), Inches(9), Inches(2))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Thank You!"
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    p2 = tf.add_paragraph()
    p2.text = "Questions & Discussion"
    p2.font.size = Pt(24)
    p2.font.color.rgb = ACCENT_BLUE
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(30)

    # Contact info
    txBox2 = slide.shapes.add_textbox(Inches(3), Inches(5.5), Inches(7), Inches(1))
    tf2 = txBox2.text_frame
    p3 = tf2.paragraphs[0]
    p3.text = "Feel free to reach out for more information"
    p3.font.size = Pt(16)
    p3.font.color.rgb = RGBColor(0xAA, 0xBB, 0xCC)
    p3.alignment = PP_ALIGN.CENTER


# Build the presentation
add_title_slide(prs)
add_agenda_slide(prs)
add_what_is_ai_slide(prs)
add_history_slide(prs)
add_types_of_ai_slide(prs)
add_applications_slide(prs)
add_ml_dl_slide(prs)
add_challenges_slide(prs)
add_future_slide(prs)
add_thank_you_slide(prs)

# Save
output_path = "/projects/sandbox/AI_Overview_Presentation.pptx"
prs.save(output_path)
print(f"Presentation saved to: {output_path}")
print(f"Total slides: {len(prs.slides)}")
